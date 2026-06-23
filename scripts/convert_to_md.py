"""
Convert the 4 source files in E:/obsidianTest/tianzhi-test/企工/ to Markdown.
- DOCX -> headings + paragraphs + tables (skip images)
- XLSX -> one Markdown section per sheet, rendered as tables
- PPTX -> per slide: heading, bullets, notes
"""
from __future__ import annotations

import os
import re
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Emu

SRC_DIR = Path(r"E:/obsidianTest/tianzhi-test/企工")

# ---------- helpers ----------

def md_escape(text: str) -> str:
    """Escape characters that would break Markdown table cells / headings."""
    if text is None:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # In headings keep text but strip newlines
    return text.strip()


def cell_to_md(value) -> str:
    if value is None:
        return ""
    s = str(value)
    # Escape pipe and collapse newlines so each cell stays a single line
    s = s.replace("|", "\\|").replace("\n", " ").replace("\r", " ")
    return s.strip()


def detect_heading_level(style_name: str | None) -> int:
    """
    Map common Word style names to Markdown heading levels.
    Heading 1 -> 1, Heading 2 -> 2, ... up to 6.
    Anything else -> 0 (regular paragraph).
    """
    if not style_name:
        return 0
    m = re.match(r"^Heading\s*(\d+)$", style_name.strip(), re.IGNORECASE)
    if m:
        return max(1, min(6, int(m.group(1))))
    return 0


# Numbered heading patterns common in Chinese enterprise docs.
# Use a list of explicit alternations rather than a single greedy pattern —
# a single pattern like ^(\d+(?:\.\d+){0,4})\s+(.+)$ exhibits catastrophic
# backtracking in Python's `re` module on inputs like "1. 标题" (P6 = None),
# while ^(\d+)\.\s+(.+)$ matches cleanly. Splitting depth-by-depth sidesteps
# the backtracking issue and is also easier to read.
HEADING_PATTERNS: list[tuple[re.Pattern, int]] = [
    (re.compile(r"^\d+\s+\S.{0,58}$"), 2),                        # 1 标题
    (re.compile(r"^\d+\.\s+\S.{0,58}$"), 2),                      # 1. 标题
    (re.compile(r"^\d+\.\d+\s+\S.{0,58}$"), 3),                   # 1.1 标题
    (re.compile(r"^\d+\.\d+\.\d+\s+\S.{0,58}$"), 4),              # 1.1.1 标题
    (re.compile(r"^\d+\.\d+\.\d+\.\d+\s+\S.{0,58}$"), 4),         # 1.1.1.1 标题
    (re.compile(r"^\d+\.\d+\.\d+\.\d+\.\d+\s+\S.{0,58}$"), 4),    # 1.1.1.1.1 标题
    # Chinese ordinals: "一、" / "二、" / "(一)" / "1、"
    (re.compile(r"^[一二三四五六七八九十百]+[、.]\s*\S.{0,58}$"), 2),
]


def pseudo_heading_level(text: str) -> int | None:
    """
    Heuristic: detect "manual" headings that the original author typed as
    plain text but rendered visually as titles (e.g. '1.1.1.1 系统简介').
    Returns Markdown heading level (2..4) or None.
    """
    s = text.strip()
    if not s or s.startswith(("http://", "https://", "www.")):
        return None
    # Reject anything that contains sentence punctuation (treating it as body)
    if re.search(r"[。，；：、]", s):
        return None
    for pat, level in HEADING_PATTERNS:
        if pat.match(s):
            return level
    return None


# ---------- DOCX ----------

def docx_to_markdown(docx_path: Path, md_path: Path) -> dict:
    doc = Document(str(docx_path))
    parts: list[str] = []
    stats = {"paragraphs": 0, "tables": 0, "headings": 0}

    # 1. Body: paragraphs and tables in document order
    body = doc.element.body
    # Build a quick lookup: para_id -> docx Paragraph object
    para_by_id = {p._p: p for p in doc.paragraphs}
    table_by_id = {t._tbl: t for t in doc.tables}

    for child in body.iterchildren():
        tag = child.tag.split("}", 1)[-1]
        if tag == "p":
            para = para_by_id.get(child)
            if para is None:
                continue
            stats["paragraphs"] += 1
            text = md_escape(para.text)
            if not text:
                parts.append("")
                continue
            level = detect_heading_level(para.style.name if para.style else None)
            if level > 0:
                stats["headings"] += 1
                parts.append(f"{'#' * level} {text}")
            else:
                # Heuristic: numbered pseudo-headings like "1.1.1.1 系统简介"
                pseudo = pseudo_heading_level(text)
                if pseudo is not None:
                    stats["headings"] += 1
                    parts.append(f"{'#' * pseudo} {text}")
                    continue
                # Honor list paragraphs
                pPr = child.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pPr")
                numPr = None
                if pPr is not None:
                    numPr = pPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}numPr")
                if numPr is not None:
                    ilvl_el = numPr.find("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ilvl")
                    ilvl = int(ilvl_el.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val", "0")) if ilvl_el is not None else 0
                    indent = "  " * ilvl
                    parts.append(f"{indent}- {text}")
                else:
                    parts.append(text)
        elif tag == "tbl":
            tbl = table_by_id.get(child)
            if tbl is None:
                continue
            stats["tables"] += 1
            rows = tbl.rows
            if not rows:
                continue
            md_rows: list[list[str]] = []
            for row in rows:
                md_rows.append([cell_to_md(cell.text) for cell in row.cells])
            if not md_rows:
                continue
            # Normalize column count
            width = max(len(r) for r in md_rows)
            for r in md_rows:
                if len(r) < width:
                    r.extend([""] * (width - len(r)))
            header = md_rows[0]
            parts.append("| " + " | ".join(header) + " |")
            parts.append("| " + " | ".join(["---"] * width) + " |")
            for r in md_rows[1:]:
                parts.append("| " + " | ".join(r) + " |")
            parts.append("")

    md_path.write_text("\n".join(parts) + "\n", encoding="utf-8")
    return stats


# ---------- XLSX ----------

def xlsx_to_markdown(xlsx_path: Path, md_path: Path) -> dict:
    wb = load_workbook(str(xlsx_path), data_only=True)
    parts: list[str] = []
    stats = {"sheets": 0, "rows": 0, "tables": 0}

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        stats["sheets"] += 1
        parts.append(f"## {sheet_name}")
        parts.append("")

        # Determine used range. We render all rows up to ws.max_row but skip
        # fully empty trailing rows.
        max_row = ws.max_row or 0
        max_col = ws.max_column or 0
        if max_row == 0 or max_col == 0:
            parts.append("_(空表)_")
            parts.append("")
            continue

        rows_md: list[list[str]] = []
        for r in range(1, max_row + 1):
            row_vals = [cell_to_md(ws.cell(row=r, column=c).value) for c in range(1, max_col + 1)]
            # If the entire row is empty, render as a blank line in the table
            rows_md.append(row_vals)
            stats["rows"] += 1

        # Trim trailing fully-empty rows
        while rows_md and all(c == "" for c in rows_md[-1]):
            rows_md.pop()
        if not rows_md:
            parts.append("_(空表)_")
            parts.append("")
            continue

        stats["tables"] += 1
        header = rows_md[0]
        parts.append("| " + " | ".join(header) + " |")
        parts.append("| " + " | ".join(["---"] * len(header)) + " |")
        for r in rows_md[1:]:
            parts.append("| " + " | ".join(r) + " |")
        parts.append("")

    md_path.write_text("\n".join(parts) + "\n", encoding="utf-8")
    return stats


# ---------- PPTX ----------

def emu_to_in(emu: int | None) -> float:
    if emu is None:
        return 0.0
    return round(emu / 914400, 2)


def pptx_to_markdown(pptx_path: Path, md_path: Path) -> dict:
    prs = Presentation(str(pptx_path))
    parts: list[str] = []
    stats = {"slides": 0, "shapes": 0, "tables": 0}

    for idx, slide in enumerate(prs.slides, 1):
        stats["slides"] += 1
        title = ""
        bullets: list[str] = []
        table_blocks: list[str] = []
        text_chunks: list[str] = []

        for shape in slide.shapes:
            stats["shapes"] += 1
            if shape.has_table:
                tbl = shape.table
                stats["tables"] += 1
                rows = list(tbl.rows)
                if not rows:
                    continue
                md_rows = [[cell_to_md(cell.text) for cell in row.cells] for row in rows]
                width = max(len(r) for r in md_rows)
                for r in md_rows:
                    if len(r) < width:
                        r.extend([""] * (width - len(r)))
                tbl_lines = [
                    "| " + " | ".join(md_rows[0]) + " |",
                    "| " + " | ".join(["---"] * width) + " |",
                ]
                for r in md_rows[1:]:
                    tbl_lines.append("| " + " | ".join(r) + " |")
                table_blocks.append("\n".join(tbl_lines))
                continue

            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            for para in tf.paragraphs:
                level = para.level or 0
                txt = md_escape(para.text)
                if not txt:
                    continue
                indent = "  " * level
                if not title and (shape == slide.shapes.title or level == 0):
                    # First top-level non-empty text often acts as the slide title
                    if not title:
                        title = txt
                        continue
                bullets.append(f"{indent}- {txt}")
                text_chunks.append(txt)

        parts.append(f"### Slide {idx} — {title or '(无标题)'}")
        if bullets:
            parts.extend(bullets)
        for tbl in table_blocks:
            parts.append("")
            parts.append(tbl)
        # Notes
        if slide.has_notes_slide:
            notes_text = md_escape(slide.notes_slide.notes_text_frame.text)
            if notes_text:
                parts.append("")
                parts.append(f"> 备注: {notes_text}")
        parts.append("")

    md_path.write_text("\n".join(parts) + "\n", encoding="utf-8")
    return stats


# ---------- driver ----------

JOBS = [
    ("企工平台建设方案_V1.2 (1).docx", docx_to_markdown, "docx"),
    ("企业级工业互联网平台产品功能清单.xlsx", xlsx_to_markdown, "xlsx"),
    ("企业级工业互联网平台产品介绍（高层版）-V26.4.pptx", pptx_to_markdown, "pptx"),
    ("企业数字化转型调研问卷-【概况信息】2025.docx", docx_to_markdown, "docx"),
]


def main() -> None:
    for fname, fn, kind in JOBS:
        src = SRC_DIR / fname
        if not src.exists():
            print(f"[SKIP] {src} not found")
            continue
        # Replace the source extension with .md; keep the rest of the stem intact
        md_name = src.stem + ".md"
        dst = SRC_DIR / md_name
        print(f"[RUN ] [{kind}] {src.name} -> {dst.name}")
        try:
            stats = fn(src, dst)
        except Exception as e:
            print(f"[FAIL] {src.name}: {e!r}")
            continue
        size_kb = round(dst.stat().st_size / 1024, 1)
        print(f"[OK  ] {dst.name}  stats={stats}  size={size_kb}KB")


if __name__ == "__main__":
    main()
